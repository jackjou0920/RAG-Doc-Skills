#!/usr/bin/env python3
"""
extract_pptx.py — Extract structured elements (title, body text, table, notes)
from a PowerPoint .pptx file using python-pptx.

Output mirrors the PDF skill's structure so both can share the same chunker and
ChromaDB collection. Each element carries: slide number (as `page`), element_type,
text, header_path (built from slide titles), and table dimensions where relevant.

Element types:
  - title : slide title (acts as a header; builds header_path)
  - body  : body/content text from text frames
  - table : table on the slide, converted to markdown
  - notes : speaker notes

Usage:
    python extract_pptx.py --input deck.pptx --output deck.json
    python extract_pptx.py --input deck.pptx          # summary only
"""

import argparse
import json
import os
import sys
from typing import Any, Dict, List


def _table_to_markdown(table) -> Dict[str, Any]:
    """Convert a python-pptx table to rows + markdown string."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append((cell.text or "").strip())
        rows.append(cells)
    if not rows:
        return {"rows": [], "markdown": "", "n_rows": 0, "n_cols": 0}
    n_cols = max(len(r) for r in rows)
    norm = [r + [""] * (n_cols - len(r)) for r in rows]
    header = norm[0]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * n_cols) + " |"]
    for r in norm[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return {"rows": norm, "markdown": "\n".join(lines), "n_rows": len(norm), "n_cols": n_cols}


def _iter_shape_text(shape) -> List[str]:
    """Collect non-empty paragraph texts from a shape's text frame."""
    texts: List[str] = []
    if not getattr(shape, "has_text_frame", False):
        return texts
    for para in shape.text_frame.paragraphs:
        runs_text = "".join(run.text for run in para.runs)
        if not runs_text and para.text:
            runs_text = para.text
        runs_text = runs_text.strip()
        if runs_text:
            texts.append(runs_text)
    return texts


def extract_pptx(pptx_path: str) -> Dict[str, Any]:
    """Extract a .pptx into a structured document dict (same shape as PDF skill)."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    elements: List[Dict[str, Any]] = []
    deck_title = os.path.splitext(os.path.basename(pptx_path))[0]

    for slide_idx, slide in enumerate(prs.slides):
        # Find slide title (if any)
        slide_title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
        except Exception:
            slide_title = ""

        header_path = f"{deck_title} > {slide_title}" if slide_title else deck_title

        if slide_title:
            elements.append(
                {
                    "page": slide_idx,
                    "element_type": "title",
                    "text": slide_title,
                    "header_path": header_path,
                }
            )

        title_shape_id = None
        try:
            if slide.shapes.title is not None:
                title_shape_id = id(slide.shapes.title)
        except Exception:
            title_shape_id = None

        # Body text + tables
        for shape in slide.shapes:
            # Skip the title shape (already captured)
            if title_shape_id is not None and id(shape) == title_shape_id:
                continue

            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = _table_to_markdown(shape.table)
                if tbl["markdown"]:
                    elements.append(
                        {
                            "page": slide_idx,
                            "element_type": "table",
                            "text": tbl["markdown"],
                            "header_path": header_path,
                            "n_rows": tbl["n_rows"],
                            "n_cols": tbl["n_cols"],
                        }
                    )
                continue

            body_texts = _iter_shape_text(shape)
            if body_texts:
                elements.append(
                    {
                        "page": slide_idx,
                        "element_type": "body",
                        "text": "\n".join(body_texts),
                        "header_path": header_path,
                    }
                )

        # Speaker notes
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    elements.append(
                        {
                            "page": slide_idx,
                            "element_type": "notes",
                            "text": notes_text,
                            "header_path": header_path,
                        }
                    )
        except Exception:
            pass

    return {
        "source": os.path.basename(pptx_path),
        "file_type": "pptx",
        "n_pages": len(prs.slides._sldIdLst) if hasattr(prs.slides, "_sldIdLst") else len(elements),
        "n_elements": len(elements),
        "elements": elements,
    }


def main():
    ap = argparse.ArgumentParser(description="Extract structured elements from a PPTX.")
    ap.add_argument("--input", required=True, help="Path to input .pptx file")
    ap.add_argument("--output", help="Path to output JSON (default: stdout summary)")
    args = ap.parse_args()

    if not os.path.isfile(args.input):
        print(f"ERROR: input not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    doc = extract_pptx(args.input)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            json.dump(doc, f, ensure_ascii=False, indent=2)
        print(f"Wrote {doc['n_elements']} elements from {doc['source']} -> {args.output}")
    else:
        counts: Dict[str, int] = {}
        for e in doc["elements"]:
            counts[e["element_type"]] = counts.get(e["element_type"], 0) + 1
        print(json.dumps({**{k: v for k, v in doc.items() if k != "elements"}, "element_counts": counts}, indent=2))


if __name__ == "__main__":
    main()